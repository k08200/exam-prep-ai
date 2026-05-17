"""
File parsing service.
Supports PDF (with image OCR fallback), PPTX, DOCX, and image files.
"""
import asyncio
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


class FileParser:
    """Extract text and metadata from lecture materials."""

    async def parse_file(self, file_path: str, file_type: str) -> dict:
        """
        Dispatch to the appropriate parser based on file_type.

        Returns:
            {"text": str, "page_count": int | None}
        """
        parsers = {
            "pdf": self.parse_pdf,
            "pptx": self.parse_pptx,
            "docx": self.parse_docx,
            "image": self.parse_image,
        }
        parser = parsers.get(file_type)
        if parser is None:
            logger.warning("No parser for file_type '%s'; returning empty text.", file_type)
            return {"text": "", "page_count": None}

        return await parser(file_path)

    async def parse_pdf(self, file_path: str) -> dict:
        """
        Extract text from a PDF using PyMuPDF.
        Falls back to pytesseract OCR for pages whose text is sparse.
        """
        def _sync_parse() -> dict:
            import fitz  # PyMuPDF

            doc = fitz.open(file_path)
            page_texts: list[str] = []

            for page in doc:
                page_text = page.get_text("text").strip()

                # If text is sparse (likely scanned page), OCR the rendered image
                if len(page_text) < 50:
                    try:
                        import pytesseract
                        from PIL import Image
                        import io

                        mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for better OCR
                        pix = page.get_pixmap(matrix=mat)
                        img_bytes = pix.tobytes("png")
                        img = Image.open(io.BytesIO(img_bytes))
                        ocr_text = pytesseract.image_to_string(img, lang="eng+kor")
                        page_texts.append(ocr_text.strip())
                    except Exception as exc:
                        logger.warning("OCR failed for PDF page: %s", exc)
                        page_texts.append(page_text)
                else:
                    page_texts.append(page_text)

            full_text = "\n\n".join(pt for pt in page_texts if pt)
            page_count = len(doc)
            doc.close()
            return {"text": full_text, "page_count": page_count}

        return await asyncio.get_event_loop().run_in_executor(None, _sync_parse)

    async def parse_pptx(self, file_path: str) -> dict:
        """
        Extract text from a PPTX file.
        Captures slide titles, body text, and speaker notes.
        """
        def _sync_parse() -> dict:
            from pptx import Presentation
            from pptx.util import Pt

            prs = Presentation(file_path)
            slide_texts: list[str] = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                parts: list[str] = [f"[Slide {slide_num}]"]

                # Title
                if slide.shapes.title and slide.shapes.title.text.strip():
                    parts.append(f"Title: {slide.shapes.title.text.strip()}")

                # All text frames
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    if shape == slide.shapes.title:
                        continue  # already captured above
                    shape_text = "\n".join(
                        para.text.strip()
                        for para in shape.text_frame.paragraphs
                        if para.text.strip()
                    )
                    if shape_text:
                        parts.append(shape_text)

                # Speaker notes
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        parts.append(f"[Notes]: {notes_text}")

                slide_texts.append("\n".join(parts))

            full_text = "\n\n".join(slide_texts)
            return {"text": full_text, "page_count": len(prs.slides)}

        return await asyncio.get_event_loop().run_in_executor(None, _sync_parse)

    async def parse_docx(self, file_path: str) -> dict:
        """
        Extract text from a DOCX file.
        Captures paragraphs and table cell content.
        """
        def _sync_parse() -> dict:
            from docx import Document

            doc = Document(file_path)
            text_parts: list[str] = []

            # Body paragraphs
            for para in doc.paragraphs:
                stripped = para.text.strip()
                if stripped:
                    text_parts.append(stripped)

            # Tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        text_parts.append(row_text)

            full_text = "\n".join(text_parts)
            # DOCX doesn't have a native "page count"; approximate by paragraph count
            return {"text": full_text, "page_count": None}

        return await asyncio.get_event_loop().run_in_executor(None, _sync_parse)

    async def parse_image(self, file_path: str) -> dict:
        """
        OCR an image file using pytesseract.
        Tries English + Korean recognition.
        """
        def _sync_parse() -> dict:
            import pytesseract
            from PIL import Image

            img = Image.open(file_path)

            # Convert to RGB if necessary (some PNGs have alpha channel)
            if img.mode not in ("RGB", "L"):
                img = img.convert("RGB")

            text = pytesseract.image_to_string(img, lang="eng+kor")
            return {"text": text.strip(), "page_count": 1}

        return await asyncio.get_event_loop().run_in_executor(None, _sync_parse)
