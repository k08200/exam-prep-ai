"""
Unit tests for AnalyticsService and FileParser service layer.
"""
import os
import tempfile

import pytest
from sqlalchemy.ext.asyncio import AsyncSession

from app.services.analytics_service import AnalyticsService
from app.services.file_parser import FileParser


@pytest.fixture
def analytics() -> AnalyticsService:
    return AnalyticsService()


@pytest.fixture
def parser() -> FileParser:
    return FileParser()


# ── AnalyticsService ───────────────────────────────────────────────────────


@pytest.mark.asyncio
async def test_update_concept_tracking_new(
    analytics: AnalyticsService, db_session: AsyncSession
) -> None:
    """First wrong answer on a concept creates a tracking record with high weakness."""
    import uuid
    user_id = str(uuid.uuid4())
    course_id = str(uuid.uuid4())

    await analytics.update_concept_tracking(
        db=db_session,
        user_id=user_id,
        course_id=course_id,
        concepts=["Gradient Descent"],
        is_correct=False,
    )
    await db_session.flush()

    from sqlalchemy import select
    from app.models.exam import ConceptTracking
    result = await db_session.execute(
        select(ConceptTracking).where(
            ConceptTracking.user_id == uuid.UUID(user_id),
            ConceptTracking.concept == "Gradient Descent",
        )
    )
    record = result.scalar_one_or_none()
    assert record is not None
    assert record.weakness_score > 0.5


@pytest.mark.asyncio
async def test_update_concept_tracking_correct_reduces_weakness(
    analytics: AnalyticsService, db_session: AsyncSession
) -> None:
    """Correct answer on a tracked concept reduces its weakness score."""
    import uuid
    user_id = str(uuid.uuid4())
    course_id = str(uuid.uuid4())

    # First: wrong answer (high weakness)
    await analytics.update_concept_tracking(
        db=db_session, user_id=user_id, course_id=course_id,
        concepts=["Backprop"], is_correct=False,
    )
    await db_session.flush()

    from sqlalchemy import select
    from app.models.exam import ConceptTracking
    result = await db_session.execute(
        select(ConceptTracking).where(
            ConceptTracking.user_id == uuid.UUID(user_id),
            ConceptTracking.concept == "Backprop",
        )
    )
    after_wrong = result.scalar_one().weakness_score

    # Then: correct answer (should reduce weakness)
    await analytics.update_concept_tracking(
        db=db_session, user_id=user_id, course_id=course_id,
        concepts=["Backprop"], is_correct=True,
    )
    await db_session.flush()

    result2 = await db_session.execute(
        select(ConceptTracking).where(
            ConceptTracking.user_id == uuid.UUID(user_id),
            ConceptTracking.concept == "Backprop",
        )
    )
    after_correct = result2.scalar_one().weakness_score
    assert after_correct < after_wrong


@pytest.mark.asyncio
async def test_get_heatmap_returns_sorted(
    analytics: AnalyticsService, db_session: AsyncSession
) -> None:
    """get_heatmap returns concepts sorted by weakness descending."""
    import uuid
    user_id = str(uuid.uuid4())
    course_id = str(uuid.uuid4())

    # Create two concepts with different weakness scores
    await analytics.update_concept_tracking(
        db=db_session, user_id=user_id, course_id=course_id,
        concepts=["Easy Topic"], is_correct=True,
    )
    await analytics.update_concept_tracking(
        db=db_session, user_id=user_id, course_id=course_id,
        concepts=["Hard Topic"], is_correct=False,
    )
    await db_session.flush()

    heatmap = await analytics.get_heatmap(db=db_session, user_id=user_id, course_id=course_id)
    assert len(heatmap) >= 1
    scores = [item["weakness_score"] for item in heatmap]
    assert scores == sorted(scores, reverse=True)


@pytest.mark.asyncio
async def test_get_heatmap_empty(
    analytics: AnalyticsService, db_session: AsyncSession
) -> None:
    """get_heatmap returns empty list when no tracking data exists."""
    import uuid
    heatmap = await analytics.get_heatmap(
        db=db_session,
        user_id=str(uuid.uuid4()),
        course_id=str(uuid.uuid4()),
    )
    assert heatmap == []


@pytest.mark.asyncio
async def test_get_cram_topics_uses_weak_concepts(analytics: AnalyticsService) -> None:
    """get_cram_topics returns topics biased toward weak user concepts."""
    analysis = {
        "top_concepts": [
            {"concept": "Gradient Descent", "importance_score": 0.9},
            {"concept": "Backpropagation", "importance_score": 0.8},
            {"concept": "Attention", "importance_score": 0.7},
        ],
        "topic_distribution": {
            "Gradient Descent": 40.0,
            "Backpropagation": 35.0,
            "Attention": 25.0,
        },
    }
    user_concepts = [
        {"concept": "Backpropagation", "weakness_score": 0.9},
        {"concept": "Gradient Descent", "weakness_score": 0.2},
    ]
    topics = await analytics.get_cram_topics(analysis=analysis, user_concepts=user_concepts)
    assert isinstance(topics, list)
    assert len(topics) > 0
    # Weakest concept (Backpropagation) should appear in cram topics
    assert "Backpropagation" in topics


@pytest.mark.asyncio
async def test_get_cram_topics_no_user_data(analytics: AnalyticsService) -> None:
    """get_cram_topics with no user data returns high-importance concepts."""
    analysis = {
        "top_concepts": [
            {"concept": "Neural Networks", "importance_score": 0.95},
        ],
        "topic_distribution": {"Neural Networks": 100.0},
    }
    topics = await analytics.get_cram_topics(analysis=analysis, user_concepts=[])
    assert isinstance(topics, list)
    assert len(topics) > 0


@pytest.mark.asyncio
async def test_get_cram_topics_intersection_priority(analytics: AnalyticsService) -> None:
    """Concepts that are both weak AND professor-important come first in cram topics."""
    analysis = {
        "top_concepts": [
            {"concept": "Gradient Descent", "importance_score": 0.95},
            {"concept": "Overfitting", "importance_score": 0.85},
        ],
        "topic_distribution": {},
    }
    user_concepts = [
        {"concept": "Gradient Descent", "weakness_score": 0.9},
        {"concept": "Regularization", "weakness_score": 0.8},
    ]
    topics = await analytics.get_cram_topics(analysis=analysis, user_concepts=user_concepts)
    # Gradient Descent is in both weak AND prof-important → should be first
    assert "Gradient Descent" in topics
    assert topics.index("Gradient Descent") == 0


def test_compute_weakness_score_zero_attempts(analytics: AnalyticsService) -> None:
    """_compute_weakness_score returns 1.0 when attempts == 0."""
    from datetime import datetime, timezone
    score = analytics._compute_weakness_score(
        attempts=0, correct_count=0,
        last_attempted=None, now=datetime.now(timezone.utc),
    )
    assert score == 1.0


def test_compute_weakness_score_timezone_naive(analytics: AnalyticsService) -> None:
    """_compute_weakness_score handles timezone-naive last_attempted."""
    from datetime import datetime, timezone
    naive_dt = datetime(2025, 1, 1, 12, 0, 0)  # no tzinfo
    now = datetime.now(timezone.utc)
    score = analytics._compute_weakness_score(
        attempts=2, correct_count=1, last_attempted=naive_dt, now=now
    )
    assert 0.0 <= score <= 1.0


def test_compute_weakness_score_all_correct(analytics: AnalyticsService) -> None:
    """Perfect score yields weakness near 0."""
    from datetime import datetime, timezone
    now = datetime.now(timezone.utc)
    score = analytics._compute_weakness_score(
        attempts=10, correct_count=10, last_attempted=now, now=now
    )
    assert score < 0.1


def test_compute_weakness_score_all_wrong(analytics: AnalyticsService) -> None:
    """All wrong yields weakness near 1."""
    from datetime import datetime, timezone
    now = datetime.now(timezone.utc)
    score = analytics._compute_weakness_score(
        attempts=10, correct_count=0, last_attempted=now, now=now
    )
    assert score > 0.7


def test_compute_weakness_score_no_last_attempted(analytics: AnalyticsService) -> None:
    """None last_attempted uses default recency weight of 1.0."""
    from datetime import datetime, timezone
    now = datetime.now(timezone.utc)
    score = analytics._compute_weakness_score(
        attempts=4, correct_count=2, last_attempted=None, now=now
    )
    assert 0.0 <= score <= 1.0


# ── FileParser ─────────────────────────────────────────────────────────────


@pytest.mark.asyncio
async def test_file_parser_pdf_text_fallback(parser: FileParser) -> None:
    """PDF parser falls back to text reading when file is plain text with .pdf ext."""
    content = "Machine learning neural networks optimization backpropagation"
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False, mode="w") as f:
        f.write(content)
        path = f.name

    try:
        result = await parser.parse_pdf(path)
        assert result["text"].strip() == content.strip()
        assert result["page_count"] == 1
    finally:
        os.unlink(path)


@pytest.mark.asyncio
async def test_file_parser_docx_with_table(parser: FileParser) -> None:
    """DOCX parser extracts table cell content."""
    import io
    from docx import Document

    doc = Document()
    doc.add_paragraph("Introduction paragraph")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Header A"
    table.cell(0, 1).text = "Header B"
    table.cell(1, 0).text = "Value 1"
    table.cell(1, 1).text = "Value 2"

    buf = io.BytesIO()
    doc.save(buf)

    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
        f.write(buf.getvalue())
        path = f.name

    try:
        result = await parser.parse_docx(path)
        assert "Introduction paragraph" in result["text"]
        assert "Header A" in result["text"]
        assert "Value 1" in result["text"]
    finally:
        os.unlink(path)


@pytest.mark.asyncio
async def test_file_parser_dispatch_pdf(parser: FileParser) -> None:
    """parse_file dispatches correctly to parse_pdf for 'pdf' type."""
    import fitz
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
        path = f.name
    try:
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Dispatch test content")
        doc.save(path)
        doc.close()

        result = await parser.parse_file(path, "pdf")
        assert "Dispatch test content" in result["text"]
        assert result["page_count"] == 1
    finally:
        os.unlink(path)


@pytest.mark.asyncio
async def test_file_parser_dispatch_unknown_type_returns_empty(parser: FileParser) -> None:
    """parse_file returns an empty result for unsupported file types."""
    result = await parser.parse_file("/tmp/not-used.xyz", "spreadsheet")

    assert result == {"text": "", "page_count": None}


@pytest.mark.asyncio
async def test_file_parser_pdf_text_fallback_failure_returns_empty(
    parser: FileParser,
) -> None:
    """PDF fallback returns an empty result when neither PDF nor text reading works."""
    result = await parser.parse_pdf("/tmp/definitely-missing-file.pdf")

    assert result == {"text": "", "page_count": None}


@pytest.mark.asyncio
async def test_file_parser_image_converts_rgba_before_ocr(
    parser: FileParser,
    monkeypatch,
) -> None:
    """Image parser converts alpha-channel images before passing them to OCR."""
    from PIL import Image

    converted_modes: list[str] = []

    def fake_ocr(image, lang: str) -> str:
        converted_modes.append(image.mode)
        assert lang == "eng+kor"
        return "ocr lecture text"

    monkeypatch.setattr("pytesseract.image_to_string", fake_ocr)

    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        path = f.name
    try:
        Image.new("RGBA", (4, 4), (255, 0, 0, 128)).save(path)

        result = await parser.parse_image(path)

        assert result == {"text": "ocr lecture text", "page_count": 1}
        assert converted_modes == ["RGB"]
    finally:
        os.unlink(path)


def test_get_claude_service_returns_mock_when_flag_set(monkeypatch) -> None:
    """get_claude_service() returns MockClaudeService when USE_MOCK_CLAUDE=True."""
    from app.core.config import settings
    from app.services import get_claude_service
    from app.services.mock_claude_service import MockClaudeService

    monkeypatch.setattr(settings, "USE_MOCK_CLAUDE", True)
    service = get_claude_service()
    assert isinstance(service, MockClaudeService)


def test_get_claude_service_requires_key_for_real_mode(monkeypatch) -> None:
    """Real Claude mode fails fast when the API key is missing."""
    from app.core.config import settings
    from app.services import get_claude_service

    monkeypatch.setattr(settings, "USE_MOCK_CLAUDE", False)
    monkeypatch.setattr(settings, "ANTHROPIC_API_KEY", "")

    with pytest.raises(RuntimeError, match="ANTHROPIC_API_KEY"):
        get_claude_service()


def test_default_claude_model_is_valid_api_id() -> None:
    """Default model uses the Anthropic API model identifier, not a display name."""
    from app.core.config import settings

    assert settings.CLAUDE_MODEL == "claude-opus-4-1-20250805"


def test_cors_origins_are_parsed_from_comma_separated_string(monkeypatch) -> None:
    """Settings exposes CORS origins as a trimmed list for FastAPI middleware."""
    from app.core.config import settings

    monkeypatch.setattr(settings, "CORS_ORIGINS", "https://app.example.com, http://localhost:3000 ")

    assert settings.cors_origins == [
        "https://app.example.com",
        "http://localhost:3000",
    ]


def test_validate_runtime_settings_allows_development_defaults(monkeypatch) -> None:
    """Development mode permits mock/local defaults."""
    from app.core.config import settings

    monkeypatch.setattr(settings, "ENVIRONMENT", "development")

    settings.validate_runtime_settings()


def test_validate_runtime_settings_rejects_production_mock_mode(monkeypatch) -> None:
    """Production mode fails fast if mock Claude is enabled."""
    from app.core.config import settings

    monkeypatch.setattr(settings, "ENVIRONMENT", "production")
    monkeypatch.setattr(settings, "SECRET_KEY", "x" * 40)
    monkeypatch.setattr(settings, "USE_MOCK_CLAUDE", True)
    monkeypatch.setattr(settings, "ANTHROPIC_API_KEY", "test-key")
    monkeypatch.setattr(settings, "AUTO_CREATE_TABLES", False)
    monkeypatch.setattr(settings, "CORS_ORIGINS", "https://app.example.com")

    with pytest.raises(RuntimeError, match="USE_MOCK_CLAUDE"):
        settings.validate_runtime_settings()


def test_validate_runtime_settings_accepts_safe_production(monkeypatch) -> None:
    """Production mode accepts explicit secret, Claude key, migrations, and CORS."""
    from app.core.config import settings

    monkeypatch.setattr(settings, "ENVIRONMENT", "production")
    monkeypatch.setattr(settings, "SECRET_KEY", "x" * 40)
    monkeypatch.setattr(settings, "USE_MOCK_CLAUDE", False)
    monkeypatch.setattr(settings, "ANTHROPIC_API_KEY", "test-key")
    monkeypatch.setattr(settings, "AUTO_CREATE_TABLES", False)
    monkeypatch.setattr(settings, "CORS_ORIGINS", "https://app.example.com")

    settings.validate_runtime_settings()


@pytest.mark.asyncio
async def test_get_current_user_rejects_token_for_missing_user(
    db_session: AsyncSession,
) -> None:
    """A valid token whose user no longer exists is rejected."""
    from fastapi import HTTPException

    from app.core.security import create_access_token, get_current_user

    token = create_access_token({"sub": "deleted-user@example.com"})

    with pytest.raises(HTTPException) as exc_info:
        await get_current_user(token=token, db=db_session)

    assert exc_info.value.status_code == 401


@pytest.mark.asyncio
async def test_get_current_user_rejects_inactive_user(
    db_session: AsyncSession,
) -> None:
    """Inactive accounts cannot authenticate even with a valid token."""
    from fastapi import HTTPException

    from app.core.security import create_access_token, get_current_user, get_password_hash
    from app.models.user import User

    user = User(
        email="inactive@example.com",
        hashed_password=get_password_hash("password123"),
        is_active=False,
    )
    db_session.add(user)
    await db_session.flush()
    token = create_access_token({"sub": user.email})

    with pytest.raises(HTTPException) as exc_info:
        await get_current_user(token=token, db=db_session)

    assert exc_info.value.status_code == 403
    assert exc_info.value.detail == "Inactive user"


def test_validate_runtime_settings_rejects_bad_timeout(monkeypatch) -> None:
    """Production mode rejects non-positive request timeout settings."""
    from app.core.config import settings

    monkeypatch.setattr(settings, "ENVIRONMENT", "production")
    monkeypatch.setattr(settings, "SECRET_KEY", "x" * 40)
    monkeypatch.setattr(settings, "USE_MOCK_CLAUDE", False)
    monkeypatch.setattr(settings, "ANTHROPIC_API_KEY", "test-key")
    monkeypatch.setattr(settings, "AUTO_CREATE_TABLES", False)
    monkeypatch.setattr(settings, "CORS_ORIGINS", "https://app.example.com")
    monkeypatch.setattr(settings, "REQUEST_TIMEOUT_SECONDS", 0)
    monkeypatch.setattr(settings, "AUTH_RATE_LIMIT_MAX_FAILURES", 5)
    monkeypatch.setattr(settings, "AUTH_RATE_LIMIT_WINDOW_SECONDS", 300)

    with pytest.raises(RuntimeError, match="REQUEST_TIMEOUT_SECONDS"):
        settings.validate_runtime_settings()


@pytest.mark.asyncio
async def test_health_endpoint(client) -> None:
    """GET /health returns liveness and AI mode metadata."""
    resp = await client.get("/health")
    assert resp.status_code == 200
    data = resp.json()
    assert data["status"] == "ok"
    assert data["ai_mode"] == "mock"
    assert data["claude_configured"] is False


@pytest.mark.asyncio
async def test_ready_endpoint_checks_dependencies(client, monkeypatch, tmp_path) -> None:
    """GET /ready checks database access and upload directory writability."""
    from app.core.config import settings

    monkeypatch.setattr(settings, "UPLOAD_DIR", str(tmp_path))

    resp = await client.get("/ready")

    assert resp.status_code == 200
    data = resp.json()
    assert data["status"] == "ready"
    assert data["database"] == "ok"
    assert data["upload_dir"] == "ok"


@pytest.mark.asyncio
async def test_root_endpoint_returns_api_metadata(client) -> None:
    """GET / returns API metadata that helps local users find docs and health checks."""
    resp = await client.get("/")

    assert resp.status_code == 200
    assert resp.json() == {
        "name": "Exam Prep AI",
        "version": "1.0.0",
        "docs": "/docs",
        "health": "/health",
    }


@pytest.mark.asyncio
async def test_health_endpoint_reports_real_claude_mode(monkeypatch) -> None:
    """health_check reflects non-mock Claude configuration."""
    from app.core.config import settings
    from app.main import health_check

    monkeypatch.setattr(settings, "USE_MOCK_CLAUDE", False)
    monkeypatch.setattr(settings, "ANTHROPIC_API_KEY", "test-key")

    data = await health_check()

    assert data["ai_mode"] == "claude"
    assert data["claude_configured"] is True


@pytest.mark.asyncio
async def test_readiness_check_reports_unwritable_upload_dir(monkeypatch, tmp_path) -> None:
    """readiness_check reports not_ready when uploads cannot be written."""
    from app.core.config import settings
    from app.main import readiness_check

    class FakeDB:
        async def execute(self, query):
            return None

    monkeypatch.setattr(settings, "UPLOAD_DIR", str(tmp_path / "missing-dir"))

    data = await readiness_check(db=FakeDB())

    assert data["status"] == "not_ready"
    assert data["database"] == "ok"
    assert data["upload_dir"] == "not_writable"


@pytest.mark.asyncio
async def test_lifespan_creates_upload_dir_runs_init_and_recovers(monkeypatch, tmp_path) -> None:
    """Application lifespan prepares local storage, migrations, and stale material recovery."""
    from unittest.mock import AsyncMock

    from fastapi import FastAPI

    from app.core.config import settings
    from app.main import lifespan

    init_mock = AsyncMock()
    recover_mock = AsyncMock(return_value=0)
    monkeypatch.setattr(settings, "UPLOAD_DIR", str(tmp_path / "uploads"))
    monkeypatch.setattr(settings, "AUTO_CREATE_TABLES", True)
    monkeypatch.setattr("app.main.init_db", init_mock)
    monkeypatch.setattr("app.main.materials.recover_stale_processing_materials", recover_mock)

    async with lifespan(FastAPI()):
        assert (tmp_path / "uploads").is_dir()

    init_mock.assert_awaited_once()
    recover_mock.assert_awaited_once()


@pytest.mark.asyncio
async def test_request_id_header_is_returned(client) -> None:
    """The request middleware returns caller-provided request IDs for tracing."""
    resp = await client.get("/health", headers={"X-Request-ID": "test-request-id"})

    assert resp.status_code == 200
    assert resp.headers["X-Request-ID"] == "test-request-id"


@pytest.mark.asyncio
async def test_request_context_middleware_times_out(monkeypatch) -> None:
    """The request middleware returns 504 when request setup exceeds timeout."""
    import asyncio
    from starlette.requests import Request
    from starlette.responses import Response

    from app.core.config import settings
    from app.core.middleware import request_context_middleware

    monkeypatch.setattr(settings, "REQUEST_TIMEOUT_SECONDS", 0.001)

    async def receive():
        return {"type": "http.request", "body": b"", "more_body": False}

    async def slow_call_next(request: Request) -> Response:
        await asyncio.sleep(0.01)
        return Response("ok")

    request = Request(
        {
            "type": "http",
            "method": "GET",
            "path": "/slow",
            "headers": [],
            "query_string": b"",
            "server": ("testserver", 80),
            "scheme": "http",
            "client": ("127.0.0.1", 1234),
        },
        receive,
    )

    response = await request_context_middleware(request, slow_call_next)

    assert response.status_code == 504
    assert response.headers["X-Request-ID"]


@pytest.mark.asyncio
async def test_request_context_middleware_success_uses_existing_request_id() -> None:
    """The request middleware preserves caller request IDs on successful responses."""
    from starlette.requests import Request
    from starlette.responses import Response

    from app.core.middleware import request_context_middleware

    async def receive():
        return {"type": "http.request", "body": b"", "more_body": False}

    async def call_next(request: Request) -> Response:
        return Response("ok", status_code=202)

    request = Request(
        {
            "type": "http",
            "method": "POST",
            "path": "/middleware-success",
            "headers": [(b"x-request-id", b"caller-request-id")],
            "query_string": b"",
            "server": ("testserver", 80),
            "scheme": "http",
            "client": ("127.0.0.1", 1234),
        },
        receive,
    )

    response = await request_context_middleware(request, call_next)

    assert response.status_code == 202
    assert response.headers["X-Request-ID"] == "caller-request-id"


@pytest.mark.asyncio
async def test_request_context_middleware_reraises_handler_errors() -> None:
    """Unexpected handler errors are logged and re-raised for FastAPI to handle."""
    from starlette.requests import Request
    from starlette.responses import Response

    from app.core.middleware import request_context_middleware

    async def receive():
        return {"type": "http.request", "body": b"", "more_body": False}

    async def failing_call_next(request: Request) -> Response:
        raise RuntimeError("handler failed")

    request = Request(
        {
            "type": "http",
            "method": "GET",
            "path": "/middleware-error",
            "headers": [],
            "query_string": b"",
            "server": ("testserver", 80),
            "scheme": "http",
            "client": ("127.0.0.1", 1234),
        },
        receive,
    )

    with pytest.raises(RuntimeError, match="handler failed"):
        await request_context_middleware(request, failing_call_next)


@pytest.mark.asyncio
async def test_sse_iter_with_heartbeat_times_out(monkeypatch) -> None:
    """SSE helper sends heartbeat events and converts stalled upstream streams to retryable errors."""
    import asyncio

    from app.core.config import settings
    from app.core.sse import iter_with_heartbeat

    monkeypatch.setattr(settings, "AI_STREAM_HEARTBEAT_SECONDS", 0.001)
    monkeypatch.setattr(settings, "AI_STREAM_EVENT_TIMEOUT_SECONDS", 0.004)

    async def stalled_events():
        await asyncio.sleep(0.02)
        yield {"type": "text", "content": "too late"}

    events = []
    async for event in iter_with_heartbeat(stalled_events()):
        events.append(event)

    assert any(event["type"] == "heartbeat" for event in events)
    assert events[-1]["type"] == "error"
    assert events[-1]["retryable"] is True


@pytest.mark.asyncio
async def test_file_parser_dispatch_pptx(parser: FileParser) -> None:
    """parse_file dispatches correctly to parse_pptx for 'pptx' type."""
    import io
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tf = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    tf.text_frame.text = "Dispatch PPTX test"
    buf = io.BytesIO()
    prs.save(buf)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        f.write(buf.getvalue())
        path = f.name

    try:
        result = await parser.parse_file(path, "pptx")
        assert "Dispatch PPTX test" in result["text"]
    finally:
        os.unlink(path)
