"""
TDD tests for material upload, listing, deletion, and file parsers.
File I/O and Claude are mocked where appropriate.
"""
import io
import os
import tempfile
import uuid
from datetime import datetime, timedelta, timezone
from pathlib import Path
from unittest.mock import AsyncMock, MagicMock, patch

import pytest
from httpx import AsyncClient
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.models.material import (
    Material,
    PROCESSING_STATUS_FAILED,
    PROCESSING_STATUS_PENDING,
    PROCESSING_STATUS_PROCESSING,
)
from app.services.file_parser import FileParser


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_tiny_pdf_bytes() -> bytes:
    """Return minimal valid PDF bytes (does not render, but parses)."""
    return (
        b"%PDF-1.4\n"
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n"
        b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] >>\nendobj\n"
        b"xref\n0 4\n0000000000 65535 f \n0000000009 00000 n \n"
        b"0000000058 00000 n \n0000000115 00000 n \n"
        b"trailer\n<< /Size 4 /Root 1 0 R >>\nstartxref\n190\n%%EOF\n"
    )


def _make_tiny_pptx_bytes() -> bytes:
    """Generate a real minimal PPTX file using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # blank
    slide = prs.slides.add_slide(slide_layout)
    tf = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    tf.text_frame.text = "Test slide content"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_tiny_docx_bytes() -> bytes:
    """Generate a real minimal DOCX file using python-docx."""
    from docx import Document

    doc = Document()
    doc.add_paragraph("Test document content")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Upload tests
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_upload_pdf_success(
    client: AsyncClient, auth_headers: dict, test_course: dict
) -> None:
    """Uploading a PDF to an owned course returns 201 with material metadata."""
    course_id = test_course["id"]
    pdf_bytes = _make_tiny_pdf_bytes()

    with patch("app.routers.materials._parse_and_update", new=AsyncMock()):
        resp = await client.post(
            f"/courses/{course_id}/materials",
            files={"files": ("lecture.pdf", pdf_bytes, "application/pdf")},
            headers=auth_headers,
        )

    assert resp.status_code == 201
    data = resp.json()
    assert len(data["materials"]) == 1
    mat = data["materials"][0]
    assert mat["original_filename"] == "lecture.pdf"
    assert mat["file_type"] == "pdf"
    assert mat["processing_status"] == "pending"
    assert mat["processing_error"] is None
    assert data["total_size"] == len(pdf_bytes)


@pytest.mark.asyncio
async def test_upload_pptx_success(
    client: AsyncClient, auth_headers: dict, test_course: dict
) -> None:
    """Uploading a PPTX file succeeds and is recognised as pptx type."""
    course_id = test_course["id"]
    pptx_bytes = _make_tiny_pptx_bytes()

    with patch("app.routers.materials._parse_and_update", new=AsyncMock()):
        resp = await client.post(
            f"/courses/{course_id}/materials",
            files={"files": ("slides.pptx", pptx_bytes, "application/vnd.ms-powerpoint")},
            headers=auth_headers,
        )

    assert resp.status_code == 201
    mat = resp.json()["materials"][0]
    assert mat["file_type"] == "pptx"
    assert mat["original_filename"] == "slides.pptx"


@pytest.mark.asyncio
async def test_upload_invalid_extension_fails(
    client: AsyncClient, auth_headers: dict, test_course: dict
) -> None:
    """Uploading a .exe file returns 422 Unprocessable Entity."""
    course_id = test_course["id"]
    resp = await client.post(
        f"/courses/{course_id}/materials",
        files={"files": ("virus.exe", b"MZ", "application/octet-stream")},
        headers=auth_headers,
    )
    assert resp.status_code == 422
    assert "not allowed" in resp.json()["detail"].lower()


@pytest.mark.asyncio
async def test_upload_too_large_fails(
    client: AsyncClient, auth_headers: dict, test_course: dict
) -> None:
    """Uploading a file exceeding 50 MB returns 413 Request Entity Too Large."""
    course_id = test_course["id"]
    # 51 MB of zeros
    big_bytes = b"\x00" * (51 * 1024 * 1024)

    resp = await client.post(
        f"/courses/{course_id}/materials",
        files={"files": ("bigfile.pdf", big_bytes, "application/pdf")},
        headers=auth_headers,
    )
    assert resp.status_code == 413


@pytest.mark.asyncio
async def test_upload_too_many_files_fails(
    client: AsyncClient, auth_headers: dict, test_course: dict, monkeypatch
) -> None:
    """Uploading more than MAX_UPLOAD_FILES returns 422 before files are parsed."""
    from app.core.config import settings

    monkeypatch.setattr(settings, "MAX_UPLOAD_FILES", 1)
    course_id = test_course["id"]

    resp = await client.post(
        f"/courses/{course_id}/materials",
        files=[
            ("files", ("one.pdf", _make_tiny_pdf_bytes(), "application/pdf")),
            ("files", ("two.pdf", _make_tiny_pdf_bytes(), "application/pdf")),
        ],
        headers=auth_headers,
    )

    assert resp.status_code == 422
    assert "at most 1 files" in resp.json()["detail"]


@pytest.mark.asyncio
async def test_upload_to_wrong_course_fails(
    client: AsyncClient, auth_headers: dict
) -> None:
    """Uploading to another user's course returns 403 Forbidden."""
    # Register a second user and create their course
    await client.post(
        "/auth/register",
        json={"email": "otheruser@example.com", "password": "otherpassword123"},
    )
    login_resp = await client.post(
        "/auth/login",
        data={"username": "otheruser@example.com", "password": "otherpassword123"},
    )
    other_token = login_resp.json()["access_token"]
    other_headers = {"Authorization": f"Bearer {other_token}"}

    course_resp = await client.post(
        "/courses",
        json={"name": "Other User's Course"},
        headers=other_headers,
    )
    other_course_id = course_resp.json()["id"]

    # Original user tries to upload to other's course
    with patch("app.routers.materials._parse_and_update", new=AsyncMock()):
        resp = await client.post(
            f"/courses/{other_course_id}/materials",
            files={"files": ("lecture.pdf", _make_tiny_pdf_bytes(), "application/pdf")},
            headers=auth_headers,
        )
    assert resp.status_code == 403


@pytest.mark.asyncio
async def test_list_materials(
    client: AsyncClient, auth_headers: dict, test_course: dict
) -> None:
    """Listing materials for a course returns all uploaded materials."""
    course_id = test_course["id"]

    with patch("app.routers.materials._parse_and_update", new=AsyncMock()):
        await client.post(
            f"/courses/{course_id}/materials",
            files={"files": ("a.pdf", _make_tiny_pdf_bytes(), "application/pdf")},
            headers=auth_headers,
        )
        await client.post(
            f"/courses/{course_id}/materials",
            files={"files": ("b.pdf", _make_tiny_pdf_bytes(), "application/pdf")},
            headers=auth_headers,
        )

    resp = await client.get(f"/courses/{course_id}/materials", headers=auth_headers)
    assert resp.status_code == 200
    materials = resp.json()
    assert len(materials) >= 2


@pytest.mark.asyncio
async def test_delete_material(
    client: AsyncClient, auth_headers: dict, test_course: dict
) -> None:
    """Deleting a material removes it from the listing."""
    course_id = test_course["id"]

    with patch("app.routers.materials._parse_and_update", new=AsyncMock()):
        upload_resp = await client.post(
            f"/courses/{course_id}/materials",
            files={"files": ("delete_me.pdf", _make_tiny_pdf_bytes(), "application/pdf")},
            headers=auth_headers,
        )
    material_id = upload_resp.json()["materials"][0]["id"]

    del_resp = await client.delete(
        f"/courses/{course_id}/materials/{material_id}",
        headers=auth_headers,
    )
    assert del_resp.status_code == 204

    list_resp = await client.get(f"/courses/{course_id}/materials", headers=auth_headers)
    ids = [m["id"] for m in list_resp.json()]
    assert material_id not in ids


@pytest.mark.asyncio
async def test_parse_failure_sets_processing_error(
    client: AsyncClient,
    auth_headers: dict,
    test_course: dict,
    db_session: AsyncSession,
) -> None:
    """Background parser failures are persisted for the UI to display."""
    from app.routers.materials import _parse_and_update

    course_id = test_course["id"]
    with patch("app.routers.materials._parse_and_update", new=AsyncMock()):
        upload_resp = await client.post(
            f"/courses/{course_id}/materials",
            files={"files": ("bad.pdf", b"not a real file", "application/pdf")},
            headers=auth_headers,
        )
    material_id = upload_resp.json()["materials"][0]["id"]
    await db_session.commit()

    class TestSessionContext:
        async def __aenter__(self):
            return db_session

        async def __aexit__(self, exc_type, exc, tb):
            return False

    with patch(
        "app.routers.materials.file_parser.parse_file",
        new=AsyncMock(side_effect=RuntimeError("parser exploded")),
    ), patch(
        "app.core.database.AsyncSessionLocal",
        return_value=TestSessionContext(),
    ):
        await _parse_and_update(uuid.UUID(material_id), "/tmp/missing.pdf", "pdf")

    result = await db_session.execute(select(Material).where(Material.id == uuid.UUID(material_id)))
    material = result.scalar_one()
    assert material.processing_status == PROCESSING_STATUS_FAILED
    assert material.processing_error == "parser exploded"


@pytest.mark.asyncio
async def test_retry_failed_material_resets_status(
    client: AsyncClient,
    auth_headers: dict,
    test_course: dict,
    db_session: AsyncSession,
) -> None:
    """Retrying a failed material clears its error and schedules parsing again."""
    course_id = test_course["id"]
    with patch("app.routers.materials._parse_and_update", new=AsyncMock()):
        upload_resp = await client.post(
            f"/courses/{course_id}/materials",
            files={"files": ("retry.pdf", _make_tiny_pdf_bytes(), "application/pdf")},
            headers=auth_headers,
        )
    material_id = upload_resp.json()["materials"][0]["id"]

    result = await db_session.execute(select(Material).where(Material.id == uuid.UUID(material_id)))
    material = result.scalar_one()
    material.processing_status = PROCESSING_STATUS_FAILED
    material.processing_error = "previous failure"
    await db_session.commit()

    with patch("app.routers.materials._parse_and_update", new=AsyncMock()) as parser_mock:
        retry_resp = await client.post(
            f"/courses/{course_id}/materials/{material_id}/retry",
            headers=auth_headers,
        )

    assert retry_resp.status_code == 200
    data = retry_resp.json()
    assert data["processing_status"] == PROCESSING_STATUS_PENDING
    assert data["processing_error"] is None
    assert parser_mock.await_count == 1


@pytest.mark.asyncio
async def test_mark_stale_processing_materials_sets_failure(
    client: AsyncClient,
    auth_headers: dict,
    test_course: dict,
    db_session: AsyncSession,
) -> None:
    """Abandoned processing materials are marked failed for user-visible retry."""
    from app.routers.materials import (
        STALE_PROCESSING_ERROR,
        mark_stale_processing_materials,
    )

    course_id = test_course["id"]
    with patch("app.routers.materials._parse_and_update", new=AsyncMock()):
        upload_resp = await client.post(
            f"/courses/{course_id}/materials",
            files={"files": ("stale.pdf", _make_tiny_pdf_bytes(), "application/pdf")},
            headers=auth_headers,
        )
    material_id = upload_resp.json()["materials"][0]["id"]

    result = await db_session.execute(select(Material).where(Material.id == uuid.UUID(material_id)))
    material = result.scalar_one()
    material.processing_status = PROCESSING_STATUS_PROCESSING
    material.created_at = datetime.now(timezone.utc) - timedelta(hours=2)
    await db_session.commit()

    marked = await mark_stale_processing_materials(db_session)

    result = await db_session.execute(select(Material).where(Material.id == uuid.UUID(material_id)))
    material = result.scalar_one()
    assert marked == 1
    assert material.processing_status == PROCESSING_STATUS_FAILED
    assert material.processing_error == STALE_PROCESSING_ERROR


# ---------------------------------------------------------------------------
# File parser unit tests
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_file_parser_pdf() -> None:
    """FileParser.parse_pdf extracts text from a real PDF on disk."""
    import fitz  # PyMuPDF

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
        path = f.name

    try:
        # Create a PDF with embedded text using PyMuPDF
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Hello World from PyMuPDF test page")
        doc.save(path)
        doc.close()

        parser = FileParser()
        result = await parser.parse_pdf(path)

        assert "text" in result
        assert "page_count" in result
        assert result["page_count"] == 1
        assert "Hello World" in result["text"]
    finally:
        os.unlink(path)


@pytest.mark.asyncio
async def test_file_parser_pptx() -> None:
    """FileParser.parse_pptx extracts text from slide content."""
    pptx_bytes = _make_tiny_pptx_bytes()

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        f.write(pptx_bytes)
        path = f.name

    try:
        parser = FileParser()
        result = await parser.parse_pptx(path)

        assert "text" in result
        assert "page_count" in result
        assert result["page_count"] == 1
        assert "Test slide content" in result["text"]
    finally:
        os.unlink(path)


@pytest.mark.asyncio
async def test_file_parser_docx() -> None:
    """FileParser.parse_docx extracts paragraph text from a DOCX."""
    docx_bytes = _make_tiny_docx_bytes()

    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
        f.write(docx_bytes)
        path = f.name

    try:
        parser = FileParser()
        result = await parser.parse_docx(path)

        assert "text" in result
        assert "Test document content" in result["text"]
    finally:
        os.unlink(path)


@pytest.mark.asyncio
async def test_file_parser_unknown_type_returns_empty() -> None:
    """Unknown file_type returns empty text without crashing."""
    parser = FileParser()
    result = await parser.parse_file("/nonexistent/file.xyz", "unknown")
    assert result["text"] == ""
    assert result["page_count"] is None
